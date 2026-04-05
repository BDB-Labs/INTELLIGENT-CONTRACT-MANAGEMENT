from __future__ import annotations

import hashlib
import logging
import subprocess
import re
import zlib
import zipfile
from dataclasses import dataclass
from pathlib import Path
from xml.etree import ElementTree

from apps.contract_intelligence.domain.enums import DocumentType
from apps.contract_intelligence.ingestion.document_classifier import classify_document

logger = logging.getLogger(__name__)

PLAIN_TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".json",
    ".yaml",
    ".yml",
    ".csv",
    ".rst",
}

# Extended format support
SPREADSHEET_SUFFIXES = {".xlsx", ".xls"}
PRESENTATION_SUFFIXES = {".pptx", ".ppt"}

IGNORED_PATH_PARTS = {
    ".git",
    ".contract_intelligence",
    ".venv",
    ".pytest_cache",
    ".ruff_cache",
    "__pycache__",
    "artifacts",
    "dist",
    "node_modules",
}


@dataclass(frozen=True)
class LoadedDocument:
    document_id: str
    relative_path: str
    document_type: DocumentType
    text: str
    text_available: bool
    text_source: str
    text_quality: str
    clauses: tuple["ClauseSpan", ...]


@dataclass(frozen=True)
class ClauseSpan:
    clause_id: str
    heading: str
    location: str
    text: str


def _document_id(relative_path: str) -> str:
    digest = hashlib.sha1(relative_path.encode("utf-8")).hexdigest()[:10]
    return f"doc_{digest}"


def _extract_docx_text(path: Path) -> str:
    with zipfile.ZipFile(path) as archive:
        raw_xml = archive.read("word/document.xml")
    root = ElementTree.fromstring(raw_xml)
    paragraphs: list[str] = []
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    for paragraph in root.findall(".//w:p", namespace):
        text_parts = [
            node.text for node in paragraph.findall(".//w:t", namespace) if node.text
        ]
        line = "".join(text_parts).strip()
        if line:
            paragraphs.append(line)
    return "\n".join(paragraphs).strip()


def _extract_xlsx_text(path: Path) -> str:
    """Extract text from Excel spreadsheets as tabular text."""
    import openpyxl

    text_parts: list[str] = []
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"\n[Sheet: {sheet_name}]")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                text_parts.append(" | ".join(c.strip() for c in cells if c.strip()))
                row_count += 1
                if row_count > 1000:
                    text_parts.append("[...truncated at 1000 rows...]")
                    break
        wb.close()
    except (OSError, IOError, ImportError) as e:
        logger.warning("XLSX extraction failed for %s: %s", path.name, e)
    return "\n".join(text_parts).strip()


def _extract_pptx_text(path: Path) -> str:
    """Extract text from PowerPoint presentations."""
    from pptx import Presentation

    text_parts: list[str] = []
    try:
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:  # type: ignore[attr-defined]
                    for para in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                        if para.text.strip():
                            text_parts.append(para.text)
    except (OSError, IOError, ImportError) as e:
        logger.warning("PPTX extraction failed for %s: %s", path.name, e)
    return "\n".join(text_parts).strip()


def _decode_pdf_literal_string(raw: str) -> str:
    raw = raw.replace(r"\(", "(").replace(r"\)", ")").replace(r"\\", "\\")
    return re.sub(r"\s+", " ", raw).strip()


def _extract_pdf_stream_text(path: Path) -> str:
    pdf_bytes = path.read_bytes()
    fragments: list[str] = []
    for match in re.finditer(
        rb"stream\r?\n(.*?)\r?\nendstream", pdf_bytes, flags=re.DOTALL
    ):
        stream = match.group(1)
        candidates = [stream]
        for wbits in (zlib.MAX_WBITS, -zlib.MAX_WBITS):
            try:
                candidates.append(zlib.decompress(stream, wbits))
            except zlib.error:
                continue
        for candidate in candidates:
            decoded = candidate.decode("latin1", errors="ignore")
            for text_match in re.finditer(r"\((?:\\.|[^\\()])*\)", decoded):
                text = _decode_pdf_literal_string(text_match.group(0)[1:-1])
                if text:
                    fragments.append(text)
    return "\n".join(fragments).strip()


def _extract_pdf_spotlight_text(path: Path) -> str:
    """Extract PDF text using macOS Spotlight metadata (macOS only)."""
    import sys

    if sys.platform != "darwin":
        logger.debug("Spotlight extraction skipped: not running on macOS")
        return ""
    try:
        result = subprocess.run(
            ["/usr/bin/mdls", "-raw", "-name", "kMDItemTextContent", str(path)],
            capture_output=True,
            text=True,
            check=False,
            timeout=30,
        )
    except OSError:
        logger.debug("Spotlight extraction failed: mdls command not found")
        return ""
    if result.returncode != 0:
        return ""
    text = result.stdout.strip()
    if not text or text == "(null)":
        return ""
    return text


def _clause_id(relative_path: str, location: str, heading: str) -> str:
    digest = hashlib.sha1(
        f"{relative_path}:{location}:{heading}".encode("utf-8")
    ).hexdigest()[:10]
    return f"clause_{digest}"


def _is_clause_heading(line: str) -> bool:
    checks = (
        r"^(section|article)\s+\d+([.\-]\d+)*[:.)]?\s+.+$",
        r"^\d+([.\-]\d+){0,4}[A-Za-z]?[:.)]?\s+.+$",
        r"^[A-Z][A-Z0-9\s/&,\-]{4,80}$",
    )
    return any(re.match(pattern, line, flags=re.IGNORECASE) for pattern in checks)


def _paragraph_clauses(relative_path: str, text: str) -> tuple[ClauseSpan, ...]:
    clauses: list[ClauseSpan] = []
    for index, block in enumerate(re.split(r"\n\s*\n", text), start=1):
        clean = "\n".join(
            line.strip() for line in block.splitlines() if line.strip()
        ).strip()
        if not clean:
            continue
        heading = clean.splitlines()[0][:80] if "\n" in clean else f"Paragraph {index}"
        location = f"{relative_path}:paragraph_{index}"
        clauses.append(
            ClauseSpan(
                clause_id=_clause_id(relative_path, location, heading),
                heading=heading,
                location=location,
                text=clean,
            )
        )
    return tuple(clauses)


def extract_clause_spans(relative_path: str, text: str) -> tuple[ClauseSpan, ...]:
    lines = [
        (line_no, line.strip())
        for line_no, line in enumerate(text.splitlines(), start=1)
        if line.strip()
    ]
    clauses: list[ClauseSpan] = []
    current_heading: str | None = None
    current_location: str | None = None
    current_lines: list[str] = []

    for line_no, line in lines:
        if _is_clause_heading(line):
            if current_heading and current_lines:
                clauses.append(
                    ClauseSpan(
                        clause_id=_clause_id(
                            relative_path,
                            current_location or f"{relative_path}:L{line_no}",
                            current_heading,
                        ),
                        heading=current_heading,
                        location=current_location or f"{relative_path}:L{line_no}",
                        text="\n".join(current_lines).strip(),
                    )
                )
            current_heading = line
            current_location = f"{relative_path}:L{line_no}"
            current_lines = [line]
        else:
            current_lines.append(line)

    if current_heading and current_lines:
        clauses.append(
            ClauseSpan(
                clause_id=_clause_id(
                    relative_path,
                    current_location or f"{relative_path}:L1",
                    current_heading,
                ),
                heading=current_heading,
                location=current_location or f"{relative_path}:L1",
                text="\n".join(current_lines).strip(),
            )
        )

    if clauses:
        return tuple(clauses)
    return _paragraph_clauses(relative_path, text)


SUPPORTED_TEXT_SUFFIXES = (
    PLAIN_TEXT_SUFFIXES
    | {".pdf", ".docx"}
    | SPREADSHEET_SUFFIXES
    | PRESENTATION_SUFFIXES
)
MAX_DOCUMENT_BYTES = 8 * 1024 * 1024
TEXT_SAMPLE_BYTES = 4096


def _is_probably_binary(path: Path) -> bool:
    sample = path.read_bytes()[:TEXT_SAMPLE_BYTES]
    if not sample:
        return False
    if b"\x00" in sample:
        return True
    control_bytes = sum(1 for byte in sample if byte < 9 or (13 < byte < 32))
    return control_bytes / max(len(sample), 1) > 0.20


def _extract_text_quality(text: str) -> str:
    normalized = re.sub(r"\s+", " ", text).strip()
    if not normalized:
        return "none"
    if len(normalized) < 20:
        return "low"
    alpha_chars = sum(1 for char in normalized if char.isalnum())
    suspicious_chars = sum(
        1
        for char in normalized
        if not (char.isalnum() or char.isspace() or char in ".,;:!?()[]{}-_/&%'\"$#@")
    )
    words = re.findall(r"[A-Za-z]{2,}", normalized)
    if len(words) < 4 or alpha_chars < 24:
        return "low"
    if suspicious_chars / max(len(normalized), 1) > 0.12:
        return "low"
    if len(words) >= 40 and len(normalized) >= 300:
        return "high"
    return "medium"


def _load_pdf_text(path: Path) -> tuple[str, str, str]:
    # Try pypdf first (proper PDF text extraction)
    try:
        import pypdf

        reader = pypdf.PdfReader(str(path))
        text_parts: list[str] = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(page_text)
        if text_parts:
            text = "\n\n".join(text_parts)
            quality = _extract_text_quality(text)
            if quality != "none":
                return text, "pypdf", quality
    except (ImportError, OSError, IOError) as exc:
        logger.debug("pypdf PDF extraction failed: %s", exc)
    except (ValueError, RuntimeError) as exc:
        logger.debug("pypdf PDF parsing error: %s", exc)

    # Fall back to existing methods
    candidates = [
        (_extract_pdf_spotlight_text(path), "pdf_spotlight"),
        (_extract_pdf_stream_text(path), "pdf_stream"),
    ]
    ranked_quality = {"high": 3, "medium": 2, "low": 1, "none": 0}
    best_text = ""
    best_source = "unavailable"
    best_quality = "none"
    for text, source in candidates:
        quality = _extract_text_quality(text)
        if ranked_quality[quality] > ranked_quality[best_quality] or (
            ranked_quality[quality] == ranked_quality[best_quality]
            and len(text) > len(best_text)
        ):
            best_text = text
            best_source = source
            best_quality = quality
    return best_text, best_source, best_quality


def _load_text(path: Path) -> tuple[str, str, str]:
    suffix = path.suffix.lower()
    if suffix in PLAIN_TEXT_SUFFIXES:
        text = path.read_text(encoding="utf-8", errors="ignore").strip()
        return text, "plain_text", _extract_text_quality(text)
    if suffix == ".docx":
        try:
            text = _extract_docx_text(path)
            return text, "docx_xml", _extract_text_quality(text)
        except (KeyError, ValueError, zipfile.BadZipFile, ElementTree.ParseError):
            return "", "unavailable", "none"
    if suffix in SPREADSHEET_SUFFIXES:
        text = _extract_xlsx_text(path)
        return text, "spreadsheet", _extract_text_quality(text)
    if suffix in PRESENTATION_SUFFIXES:
        text = _extract_pptx_text(path)
        return text, "presentation", _extract_text_quality(text)
    if suffix == ".pdf":
        return _load_pdf_text(path)
    if not _is_probably_binary(path):
        text = path.read_text(encoding="utf-8", errors="ignore").strip()
        return text, "plain_text_fallback", _extract_text_quality(text)
    return "", "unavailable", "none"


def iter_project_documents(project_dir: Path) -> list[LoadedDocument]:
    documents: list[LoadedDocument] = []
    for path in sorted(project_dir.rglob("*")):
        if not path.is_file():
            continue
        relative_path = path.relative_to(project_dir).as_posix()
        if any(
            part in IGNORED_PATH_PARTS for part in path.relative_to(project_dir).parts
        ):
            continue
        suffix = path.suffix.lower()
        classified_from_name = classify_document(path.name)
        if path.stat().st_size > MAX_DOCUMENT_BYTES:
            if (
                suffix not in SUPPORTED_TEXT_SUFFIXES
                and classified_from_name is DocumentType.OTHER
            ):
                continue
            text, text_source, text_quality = "", "skipped_oversize", "none"
        elif suffix not in SUPPORTED_TEXT_SUFFIXES and _is_probably_binary(path):
            if classified_from_name is DocumentType.OTHER:
                continue
            text, text_source, text_quality = "", "unsupported_binary", "none"
        else:
            text, text_source, text_quality = _load_text(path)
        text_available = bool(text) and text_quality != "low"
        clauses = extract_clause_spans(relative_path, text) if text_available else ()
        document_type = classify_document(path.name, text=text)
        documents.append(
            LoadedDocument(
                document_id=_document_id(relative_path),
                relative_path=relative_path,
                document_type=document_type,
                text=text,
                text_available=text_available,
                text_source=text_source if text else text_source,
                text_quality=text_quality,
                clauses=clauses,
            )
        )
    return documents
