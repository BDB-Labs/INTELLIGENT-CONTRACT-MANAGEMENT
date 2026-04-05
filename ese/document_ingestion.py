"""Document ingestion pipeline with multi-format support and self-correction."""

from __future__ import annotations

import json
import logging
import re
import textwrap
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

from ese.adapters import AdapterExecutionError

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".doc",
    ".xlsx",
    ".xls",
    ".pptx",
    ".ppt",
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".rtf",
    ".eml",
    ".msg",
}

TEXT_EXTRACTABLE = {".pdf", ".docx", ".doc", ".rtf"}
SPREADSHEET_FORMATS = {".xlsx", ".xls", ".csv"}
PRESENTATION_FORMATS = {".pptx", ".ppt"}
MARKUP_FORMATS = {".html", ".htm", ".xml", ".md", ".json"}
PLAIN_TEXT_FORMATS = {".txt", ".eml", ".msg"}


@dataclass(frozen=True)
class ExtractedDocument:
    """Result of document extraction."""

    source_path: Path
    filename: str
    extension: str
    text_content: str
    metadata: dict[str, Any] = field(default_factory=dict)
    extraction_quality: str = "good"
    warnings: list[str] = field(default_factory=list)
    clause_boundaries: list[tuple[int, int, str]] = field(default_factory=list)


@dataclass(frozen=True)
class ReflectionResult:
    """Result of self-reflection on model output."""

    role: str
    original_output: str
    reflection_score: float
    issues_found: list[str]
    improvement_suggestions: list[str]
    needs_regeneration: bool
    refined_output: str | None = None


def extract_text_from_pdf(path: Path) -> tuple[str, list[str]]:
    """Extract text from PDF with fallback strategies."""
    import pypdf

    warnings: list[str] = []
    text_parts: list[str] = []

    try:
        reader = pypdf.PdfReader(str(path))
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(page_text)
                else:
                    warnings.append(
                        f"Page {i + 1}: No extractable text (may be scanned/image-based)"
                    )
            except (OSError, IOError, ValueError) as e:
                warnings.append(f"Page {i + 1}: Extraction error - {e}")

        if not text_parts:
            warnings.append(
                "PDF contains no extractable text content. Consider OCR preprocessing."
            )
    except (OSError, IOError) as e:
        warnings.append(f"PDF read error: {e}")

    return "\n\n".join(text_parts), warnings


def extract_text_from_docx(path: Path) -> tuple[str, list[str]]:
    """Extract text from Word document preserving structure."""
    from docx import Document

    warnings: list[str] = []
    text_parts: list[str] = []

    try:
        doc = Document(str(path))
        for para in doc.paragraphs:
            if para.text.strip():
                if para.style.name.startswith("Heading"):
                    text_parts.append(f"\n{'#' * int(para.style.name[-1])} {para.text}")
                else:
                    text_parts.append(para.text)

        for table in doc.tables:
            text_parts.append("\n[Table]")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                text_parts.append(" | ".join(cells))
            text_parts.append("[/Table]")

        if not text_parts:
            warnings.append("Document appears empty or contains only images")
    except (OSError, IOError, ImportError) as e:
        warnings.append(f"DOCX read error: {e}")

    return "\n".join(text_parts), warnings


def extract_text_from_spreadsheet(path: Path) -> tuple[str, list[str]]:
    """Extract structured text from spreadsheets."""
    import csv

    warnings: list[str] = []
    text_parts: list[str] = []
    ext = path.suffix.lower()

    try:
        if ext == ".csv":
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    text_parts.append(" | ".join(cell.strip() for cell in row))
                    if i > 500:
                        warnings.append(
                            "Spreadsheet truncated at 500 rows for extraction"
                        )
                        break
        else:
            import openpyxl

            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"\n[Sheet: {sheet_name}]")
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    text_parts.append(" | ".join(c.strip() for c in cells if c.strip()))
                    row_count += 1
                    if row_count > 500:
                        warnings.append(f"Sheet '{sheet_name}' truncated at 500 rows")
                        break
            wb.close()

        if not text_parts:
            warnings.append("Spreadsheet appears empty")
    except (OSError, IOError, ImportError) as e:
        warnings.append(f"Spreadsheet read error: {e}")

    return "\n".join(text_parts), warnings


def extract_text_from_presentation(path: Path) -> tuple[str, list[str]]:
    """Extract text from presentations."""
    from pptx import Presentation

    warnings: list[str] = []
    text_parts: list[str] = []

    try:
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)

        if not text_parts:
            warnings.append("Presentation appears empty or contains only images")
    except (OSError, IOError, ImportError) as e:
        warnings.append(f"Presentation read error: {e}")

    return "\n".join(text_parts), warnings


def extract_text_from_markup(path: Path) -> tuple[str, list[str]]:
    """Extract text from markup formats."""
    warnings: list[str] = []
    ext = path.suffix.lower()

    try:
        content = path.read_text(encoding="utf-8", errors="replace")

        if ext in {".html", ".htm"}:
            import re

            content = re.sub(r"<script[^>]*>.*?</script>", "", content, flags=re.DOTALL)
            content = re.sub(r"<style[^>]*>.*?</style>", "", content, flags=re.DOTALL)
            content = re.sub(r"<[^>]+>", " ", content)
            content = re.sub(r"\s+", " ", content).strip()
        elif ext == ".xml":
            import re

            content = re.sub(r"<[^>]+>", " ", content)
            content = re.sub(r"\s+", " ", content).strip()

        if not content.strip():
            warnings.append("File appears empty after extraction")

        return content, warnings
    except (OSError, IOError) as e:
        warnings.append(f"Markup read error: {e}")
        return "", warnings


def extract_text_from_plain(path: Path) -> tuple[str, list[str]]:
    """Read plain text files directly."""
    try:
        content = path.read_text(encoding="utf-8", errors="replace")
        return content, []
    except (OSError, IOError) as e:
        return "", [f"Plain text read error: {e}"]


EXTRACTORS: dict[str, Callable[[Path], tuple[str, list[str]]]] = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".doc": extract_text_from_docx,
    ".rtf": extract_text_from_plain,
    ".xlsx": extract_text_from_spreadsheet,
    ".xls": extract_text_from_spreadsheet,
    ".csv": extract_text_from_spreadsheet,
    ".pptx": extract_text_from_presentation,
    ".ppt": extract_text_from_presentation,
    ".html": extract_text_from_markup,
    ".htm": extract_text_from_markup,
    ".xml": extract_text_from_markup,
    ".md": extract_text_from_plain,
    ".json": extract_text_from_plain,
    ".txt": extract_text_from_plain,
    ".eml": extract_text_from_plain,
    ".msg": extract_text_from_plain,
}


def extract_document(path: Path, *, max_chars: int = 500_000) -> ExtractedDocument:
    """Extract text from a document file with format detection."""
    ext = path.suffix.lower()
    filename = path.name

    if ext not in EXTRACTORS:
        return ExtractedDocument(
            source_path=path,
            filename=filename,
            extension=ext,
            text_content="",
            extraction_quality="unsupported",
            warnings=[f"Unsupported file format: {ext}"],
        )

    extractor = EXTRACTORS[ext]
    text_content, warnings = extractor(path)

    if len(text_content) > max_chars:
        text_content = text_content[:max_chars]
        warnings.append(f"Content truncated to {max_chars} characters")

    quality = "good"
    if not text_content.strip():
        quality = "empty"
    elif len(warnings) > 2:
        quality = "poor"
    elif warnings:
        quality = "fair"

    return ExtractedDocument(
        source_path=path,
        filename=filename,
        extension=ext,
        text_content=text_content,
        extraction_quality=quality,
        warnings=warnings,
    )


def detect_clause_boundaries(text: str) -> list[tuple[int, int, str]]:
    """Detect clause boundaries in contract text."""
    boundaries: list[tuple[int, int, str]] = []

    heading_patterns = [
        re.compile(r"^(Article\s+[IVXLCDM]+[.:]\s*.+)$", re.MULTILINE | re.IGNORECASE),
        re.compile(r"^(Section\s+\d+[.:]\s*.+)$", re.MULTILINE | re.IGNORECASE),
        re.compile(r"^(Section\s+[A-Z][.:]\s*.+)$", re.MULTILINE | re.IGNORECASE),
        re.compile(r"^(\d+\.\s+[A-Z].+)$", re.MULTILINE),
        re.compile(r"^([A-Z][A-Z\s]{3,})$", re.MULTILINE),
    ]

    matches: list[tuple[int, str]] = []
    for pattern in heading_patterns:
        for match in pattern.finditer(text):
            matches.append((match.start(), match.group(1).strip()))

    matches.sort(key=lambda x: x[0])

    for i, (start, heading) in enumerate(matches):
        end = matches[i + 1][0] if i + 1 < len(matches) else len(text)
        boundaries.append((start, end, heading))

    return boundaries


def chunk_document(
    text: str, *, max_chunk_size: int = 4000, overlap: int = 200
) -> list[str]:
    """Split document into overlapping chunks for processing."""
    if len(text) <= max_chunk_size:
        return [text]

    chunks: list[str] = []
    start = 0

    while start < len(text):
        end = start + max_chunk_size

        if end >= len(text):
            chunks.append(text[start:])
            break

        break_point = text.rfind("\n\n", start + max_chunk_size - overlap, end)
        if break_point == -1:
            break_point = text.rfind("\n", start + max_chunk_size - overlap, end)
        if break_point == -1:
            break_point = text.rfind(". ", start + max_chunk_size - overlap, end)
            if break_point != -1:
                break_point += 2

        if break_point == -1 or break_point <= start:
            break_point = end

        chunks.append(text[start:break_point].strip())
        start = break_point - overlap if break_point > overlap else break_point

    return chunks


def reflect_on_output(
    *,
    role: str,
    output: str,
    role_prompt: str,
    adapter: Callable[..., str],
    model: str,
    cfg: dict[str, Any],
    context: dict[str, str] | None = None,
    max_reflection_rounds: int = 2,
    min_acceptable_score: float = 0.6,
) -> ReflectionResult:
    """Self-reflect on model output and suggest improvements.

    Uses the same adapter to have the model critique its own work.
    If the score is below threshold, regenerates with improvement suggestions.
    """
    reflection_prompt = textwrap.dedent(
        f"""
        You are reviewing your own previous output for quality and completeness.

        ROLE: {role}
        ORIGINAL PROMPT: {role_prompt[:2000]}

        YOUR PREVIOUS OUTPUT:
        {output[:8000]}

        Evaluate your output on these criteria (score 0.0 to 1.0):
        1. COMPLETENESS: Did you address all aspects of the prompt?
        2. EVIDENCE: Are your claims supported by the provided context?
        3. SPECIFICITY: Is your output specific and actionable, not generic?
        4. STRUCTURE: Does your output follow the expected format?
        5. CERTAINTY: Did you properly flag uncertainty where evidence is thin?

        Return a JSON object with:
        - overall_score: float 0.0-1.0
        - issues: list of specific issues found
        - suggestions: list of concrete improvements
        - needs_regeneration: boolean (true if score < 0.6)
    """
    ).strip()

    score = 0.5
    issues: list[str] = []
    suggestions: list[str] = []

    for round_num in range(max_reflection_rounds + 1):
        try:
            reflection_text = adapter(
                role=f"{role}_self_review",
                model=model,
                prompt=reflection_prompt,
                context=context or {},
                cfg=cfg,
            )

            reflection_data = _parse_reflection_json(reflection_text)
            if reflection_data is None:
                logger.warning(
                    "Reflection output was not valid JSON, skipping self-correction"
                )
                return ReflectionResult(
                    role=role,
                    original_output=output,
                    reflection_score=0.5,
                    issues_found=["Self-reflection output was not parseable"],
                    improvement_suggestions=[],
                    needs_regeneration=False,
                )

            score = reflection_data.get("overall_score", 0.5)
            issues = reflection_data.get("issues", [])
            suggestions = reflection_data.get("suggestions", [])
            needs_regeneration = reflection_data.get("needs_regeneration", False)

            logger.info(
                "Reflection round %d for role=%s: score=%.2f issues=%d needs_regeneration=%s",
                round_num + 1,
                role,
                score,
                len(issues),
                needs_regeneration,
            )

            if score >= min_acceptable_score or round_num >= max_reflection_rounds:
                return ReflectionResult(
                    role=role,
                    original_output=output,
                    reflection_score=score,
                    issues_found=issues,
                    improvement_suggestions=suggestions,
                    needs_regeneration=needs_regeneration
                    and round_num < max_reflection_rounds,
                    refined_output=output if not needs_regeneration else None,
                )

            regeneration_prompt = textwrap.dedent(
                f"""
                Your previous output was reviewed and scored {score}/1.0.
                Issues found: {json.dumps(issues)}
                Suggested improvements: {json.dumps(suggestions)}

                ORIGINAL PROMPT: {role_prompt[:2000]}

                Regenerate your output addressing ALL the issues above.
                Be more specific, evidence-based, and complete.
                Flag uncertainty where evidence is thin.
            """
            ).strip()

            output = adapter(
                role=role,
                model=model,
                prompt=regeneration_prompt,
                context=context or {},
                cfg=cfg,
            )

        except (AdapterExecutionError, json.JSONDecodeError, TimeoutError) as e:
            logger.warning("Reflection failed for role=%s: %s", role, e)
            return ReflectionResult(
                role=role,
                original_output=output,
                reflection_score=0.5,
                issues_found=[f"Reflection process failed: {e}"],
                improvement_suggestions=[],
                needs_regeneration=False,
            )

    return ReflectionResult(
        role=role,
        original_output=output,
        reflection_score=score,
        issues_found=issues,
        improvement_suggestions=suggestions,
        needs_regeneration=False,
        refined_output=output,
    )


def _parse_reflection_json(text: str) -> dict[str, Any] | None:
    """Extract JSON from reflection output."""
    import json

    text = text.strip()

    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    json_match = re.search(r"\{[\s\S]*\}", text)
    if json_match:
        try:
            return json.loads(json_match.group())
        except json.JSONDecodeError:
            pass

    return None


def ingest_documents(
    directory: str | Path,
    *,
    max_chars_per_doc: int = 500_000,
    extensions: set[str] | None = None,
) -> list[ExtractedDocument]:
    """Ingest all supported documents from a directory."""
    dir_path = Path(directory)
    if not dir_path.is_dir():
        raise ValueError(f"Not a directory: {directory}")

    allowed = extensions or SUPPORTED_EXTENSIONS
    documents: list[ExtractedDocument] = []

    for path in sorted(dir_path.iterdir()):
        if not path.is_file():
            continue
        if path.suffix.lower() not in allowed:
            continue

        doc = extract_document(path, max_chars=max_chars_per_doc)
        documents.append(doc)

        if doc.warnings:
            for warning in doc.warnings:
                logger.warning("Document %s: %s", doc.filename, warning)

        logger.info(
            "Extracted %s: %d chars, quality=%s, warnings=%d",
            doc.filename,
            len(doc.text_content),
            doc.extraction_quality,
            len(doc.warnings),
        )

    return documents
